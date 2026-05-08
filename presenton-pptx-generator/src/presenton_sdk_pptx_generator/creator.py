import os
import shutil
import tempfile
import uuid
import zipfile
from typing import List, Optional

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.text.text import Font, TextFrame, _Paragraph, _Run
from pptx.util import Pt

from .download import (
    download_files,
    is_data_image_uri,
    is_http_url,
    resolve_image_path_to_filesystem,
    write_data_image_to_file,
)
from .image_utils import (
    clip_image,
    create_circle_image,
    fit_image,
    invert_image,
    round_image_corners,
    set_image_opacity,
)
from .models import (
    PptxAutoShapeBoxModel,
    PptxBoxShapeEnum,
    PptxConnectorModel,
    PptxFillModel,
    PptxFontModel,
    PptxParagraphModel,
    PptxPictureBoxModel,
    PptxPositionModel,
    PptxPresentationModel,
    PptxShadowModel,
    PptxSlideModel,
    PptxSpacingModel,
    PptxStrokeModel,
    PptxTextBoxModel,
    PptxTextRunModel,
)
from .text_runs import parse_html_text_to_text_runs

BLANK_SLIDE_LAYOUT = 6


class PptxPresentationCreator:
    def __init__(self, ppt_model: PptxPresentationModel, temp_dir: str):
        self._temp_dir = temp_dir
        self._ppt_model = ppt_model
        self._slide_models = ppt_model.slides

        self._ppt = Presentation()
        self._ppt.slide_width = Pt(1280)
        self._ppt.slide_height = Pt(720)

    def get_sub_element(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def fix_keynote_compatibility(self, pptx_path: str):
        presentation_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        rel_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        package_rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        content_types_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
        notes_master_rel_type = (
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster"
        )
        printer_settings_rel_type = (
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/printerSettings"
        )
        printer_settings_content_type = (
            "application/vnd.openxmlformats-officedocument.presentationml.printerSettings"
        )

        def ensure_grp_sppr_xfrm(slide_path: str):
            slide_tree = etree.parse(slide_path)
            slide_root = slide_tree.getroot()
            grp_sppr_elements = slide_root.findall(f".//{{{presentation_ns}}}grpSpPr")
            changed = False
            for grp_sppr in grp_sppr_elements:
                xfrm = grp_sppr.find(f"{{{drawing_ns}}}xfrm")
                if xfrm is None:
                    xfrm = etree.SubElement(grp_sppr, f"{{{drawing_ns}}}xfrm")
                    etree.SubElement(xfrm, f"{{{drawing_ns}}}off", x="0", y="0")
                    etree.SubElement(xfrm, f"{{{drawing_ns}}}ext", cx="0", cy="0")
                    etree.SubElement(xfrm, f"{{{drawing_ns}}}chOff", x="0", y="0")
                    etree.SubElement(xfrm, f"{{{drawing_ns}}}chExt", cx="0", cy="0")
                    changed = True
            if changed:
                slide_tree.write(
                    slide_path,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone="yes",
                )

        def ensure_drawing_paragraph_end_rpr(xml_path: str):
            xml_tree = etree.parse(xml_path)
            xml_root = xml_tree.getroot()
            paragraph_elements = xml_root.findall(f".//{{{drawing_ns}}}p")
            changed = False
            end_para_rpr_tag = f"{{{drawing_ns}}}endParaRPr"

            for paragraph in paragraph_elements:
                end_para_rpr_elements = paragraph.findall(end_para_rpr_tag)
                if not end_para_rpr_elements:
                    etree.SubElement(paragraph, end_para_rpr_tag)
                    changed = True
                    continue

                keep_end_para_rpr = end_para_rpr_elements[0]
                for duplicate_end_para_rpr in end_para_rpr_elements[1:]:
                    paragraph.remove(duplicate_end_para_rpr)
                    changed = True

                if len(paragraph) == 0 or paragraph[-1] is keep_end_para_rpr:
                    continue

                paragraph.remove(keep_end_para_rpr)
                paragraph.append(keep_end_para_rpr)
                changed = True

            if changed:
                xml_tree.write(
                    xml_path,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone="yes",
                )

        def ensure_slide_paragraphs_end_rpr(ppt_dir: str):
            slides_dir = os.path.join(ppt_dir, "slides")
            if not os.path.isdir(slides_dir):
                return

            for file_name in os.listdir(slides_dir):
                if file_name.endswith(".xml"):
                    ensure_drawing_paragraph_end_rpr(os.path.join(slides_dir, file_name))

        def normalize_zero_shadows(xml_path: str):
            xml_tree = etree.parse(xml_path)
            xml_root = xml_tree.getroot()
            changed = False

            for outer_shadow in xml_root.findall(f".//{{{drawing_ns}}}outerShdw"):
                if not all(
                    outer_shadow.get(attr_name) == "0"
                    for attr_name in ("blurRad", "dist", "dir")
                ):
                    continue

                color_element = outer_shadow.find(f"{{{drawing_ns}}}srgbClr")
                alpha_element = (
                    color_element.find(f"{{{drawing_ns}}}alpha")
                    if color_element is not None
                    else None
                )
                zero_shadow_attrs = ["blurRad", "dist", "dir"]
                if alpha_element is not None and alpha_element.get("val") == "0":
                    zero_shadow_attrs.append("rotWithShape")

                for attr_name in zero_shadow_attrs:
                    if attr_name in outer_shadow.attrib:
                        del outer_shadow.attrib[attr_name]
                        changed = True

            if changed:
                xml_tree.write(
                    xml_path,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone="yes",
                )

        def normalize_zero_shadows_in_pptx(extract_dir: str):
            ppt_dir = os.path.join(extract_dir, "ppt")
            if not os.path.isdir(ppt_dir):
                return

            for root, _, files in os.walk(ppt_dir):
                for file_name in files:
                    if not file_name.endswith(".xml"):
                        continue
                    normalize_zero_shadows(os.path.join(root, file_name))

        def theme_is_complete(theme_path: str):
            theme_tree = etree.parse(theme_path)
            theme_root = theme_tree.getroot()
            required_lists = (
                "fillStyleLst",
                "lnStyleLst",
                "effectStyleLst",
                "bgFillStyleLst",
            )
            for list_name in required_lists:
                style_list = theme_root.find(f".//{{{drawing_ns}}}{list_name}")
                if style_list is None or len(style_list) < 3:
                    return False

            major_font = theme_root.find(f".//{{{drawing_ns}}}majorFont")
            minor_font = theme_root.find(f".//{{{drawing_ns}}}minorFont")
            if major_font is None or minor_font is None:
                return False
            if (
                len(major_font.findall(f"{{{drawing_ns}}}font")) < 10
                or len(minor_font.findall(f"{{{drawing_ns}}}font")) < 10
            ):
                return False

            return True

        def normalize_incomplete_themes(ppt_dir: str):
            themes_dir = os.path.join(ppt_dir, "theme")
            if not os.path.isdir(themes_dir):
                return

            theme_paths = [
                os.path.join(themes_dir, file_name)
                for file_name in os.listdir(themes_dir)
                if file_name.startswith("theme") and file_name.endswith(".xml")
            ]
            complete_theme_path = next(
                (theme_path for theme_path in theme_paths if theme_is_complete(theme_path)),
                None,
            )
            if complete_theme_path is None:
                return

            for theme_path in theme_paths:
                if theme_path == complete_theme_path or theme_is_complete(theme_path):
                    continue
                shutil.copyfile(complete_theme_path, theme_path)

        def remove_printer_settings_parts(ppt_dir: str, targets: List[str]):
            for target in targets:
                target_path = os.path.normpath(
                    os.path.join(ppt_dir, target.replace("/", os.sep))
                )
                try:
                    common_path = os.path.commonpath([ppt_dir, target_path])
                except ValueError:
                    continue
                if common_path != ppt_dir:
                    continue
                if os.path.isfile(target_path):
                    os.remove(target_path)
                    parent = os.path.dirname(target_path)
                    while parent != ppt_dir and os.path.isdir(parent):
                        try:
                            os.rmdir(parent)
                        except OSError:
                            break
                        parent = os.path.dirname(parent)

            printer_settings_dir = os.path.join(ppt_dir, "printerSettings")
            if os.path.isdir(printer_settings_dir):
                shutil.rmtree(printer_settings_dir)

        def remove_printer_settings_content_type(extract_dir: str):
            has_bin_parts = False
            for root, _, files in os.walk(extract_dir):
                if any(file_name.lower().endswith(".bin") for file_name in files):
                    has_bin_parts = True
                    break

            if has_bin_parts:
                return

            content_types_path = os.path.join(extract_dir, "[Content_Types].xml")
            if not os.path.exists(content_types_path):
                return

            content_types_tree = etree.parse(content_types_path)
            content_types_root = content_types_tree.getroot()
            default_tag = f"{{{content_types_ns}}}Default"
            changed = False
            for default in list(content_types_root.findall(default_tag)):
                if (
                    default.get("Extension") == "bin"
                    and default.get("ContentType") == printer_settings_content_type
                ):
                    content_types_root.remove(default)
                    changed = True

            if changed:
                content_types_tree.write(
                    content_types_path,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone="yes",
                )

        def normalize_presentation_relationship_ids(
            rels_path: str,
            presentation_path: str,
        ):
            rels_tree = etree.parse(rels_path)
            rels_root = rels_tree.getroot()
            presentation_tree = etree.parse(presentation_path)
            presentation_root = presentation_tree.getroot()
            rel_tag = f"{{{package_rel_ns}}}Relationship"
            rel_id_attr = f"{{{rel_ns}}}id"

            rels = list(rels_root.findall(rel_tag))
            original_relationships = [
                (rel.get("Id"), rel.get("Type"), rel.get("Target")) for rel in rels
            ]
            rels_by_id = {
                rel.get("Id"): rel for rel in rels if rel.get("Id") is not None
            }
            assigned_ids = {}
            ordered_rels = []
            next_id = 1

            def assign_rel(rel, preferred_id: Optional[str] = None):
                nonlocal next_id
                if rel is None or rel in ordered_rels:
                    return None

                rel_id = preferred_id
                if rel_id is None:
                    while f"rId{next_id}" in assigned_ids:
                        next_id += 1
                    rel_id = f"rId{next_id}"

                rel.set("Id", rel_id)
                assigned_ids[rel_id] = rel
                ordered_rels.append(rel)

                if rel_id.startswith("rId") and rel_id[3:].isdigit():
                    next_id = max(next_id, int(rel_id[3:]) + 1)

                return rel_id

            slide_master_ids = presentation_root.findall(
                f".//{{{presentation_ns}}}sldMasterId"
            )
            for slide_master_id in slide_master_ids:
                rel = rels_by_id.get(slide_master_id.get(rel_id_attr))
                new_rel_id = assign_rel(rel, f"rId{next_id}")
                if new_rel_id:
                    slide_master_id.set(rel_id_attr, new_rel_id)

            slide_ids = presentation_root.findall(f".//{{{presentation_ns}}}sldId")
            for slide_id in slide_ids:
                rel = rels_by_id.get(slide_id.get(rel_id_attr))
                new_rel_id = assign_rel(rel)
                if new_rel_id:
                    slide_id.set(rel_id_attr, new_rel_id)

            notes_master_ids = presentation_root.findall(
                f".//{{{presentation_ns}}}notesMasterId"
            )
            for notes_master_id in notes_master_ids:
                rel = rels_by_id.get(notes_master_id.get(rel_id_attr))
                new_rel_id = assign_rel(rel)
                if new_rel_id:
                    notes_master_id.set(rel_id_attr, new_rel_id)

            known_tail_types = [
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/presProps",
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/viewProps",
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme",
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tableStyles",
            ]
            for rel_type in known_tail_types:
                for rel in rels:
                    if rel.get("Type") == rel_type:
                        assign_rel(rel)

            for rel in rels:
                assign_rel(rel)

            normalized_relationships = [
                (rel.get("Id"), rel.get("Type"), rel.get("Target"))
                for rel in ordered_rels
            ]
            if normalized_relationships != original_relationships:
                for rel in list(rels_root.findall(rel_tag)):
                    rels_root.remove(rel)
                for rel in ordered_rels:
                    rels_root.append(rel)
                rels_tree.write(
                    rels_path,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone="yes",
                )
                presentation_tree.write(
                    presentation_path,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone="yes",
                )

        with tempfile.TemporaryDirectory() as temp_dir:
            extract_dir = os.path.join(temp_dir, "pptx_contents")
            os.makedirs(extract_dir, exist_ok=True)
            with zipfile.ZipFile(pptx_path, "r") as existing_zip:
                existing_zip.extractall(extract_dir)

            ppt_dir = os.path.join(extract_dir, "ppt")
            slides_dir = os.path.join(ppt_dir, "slides")
            if os.path.isdir(slides_dir):
                for file_name in os.listdir(slides_dir):
                    if file_name.endswith(".xml"):
                        ensure_grp_sppr_xfrm(os.path.join(slides_dir, file_name))

            rels_path = os.path.join(ppt_dir, "_rels", "presentation.xml.rels")
            presentation_path = os.path.join(ppt_dir, "presentation.xml")
            if os.path.exists(rels_path) and os.path.exists(presentation_path):
                rels_tree = etree.parse(rels_path)
                rels_root = rels_tree.getroot()
                rel_tag = f"{{{package_rel_ns}}}Relationship"
                notes_master_rel = None
                existing_ids = set()
                rels_changed = False
                printer_settings_targets = []
                for rel in list(rels_root.findall(rel_tag)):
                    rel_id = rel.get("Id")
                    if rel_id:
                        existing_ids.add(rel_id)
                    if rel.get("Type") == notes_master_rel_type:
                        notes_master_rel = rel
                    if rel.get("Type") == printer_settings_rel_type:
                        target = rel.get("Target")
                        if target:
                            printer_settings_targets.append(target)
                        rels_root.remove(rel)
                        rels_changed = True

                notes_masters_dir = os.path.join(ppt_dir, "notesMasters")
                has_notes_master = os.path.isdir(notes_masters_dir) and any(
                    name.endswith(".xml") for name in os.listdir(notes_masters_dir)
                )

                if has_notes_master and notes_master_rel is None:
                    next_id = 1
                    while f"rId{next_id}" in existing_ids:
                        next_id += 1
                    notes_master_rel = etree.SubElement(rels_root, rel_tag)
                    notes_master_rel.set("Id", f"rId{next_id}")
                    notes_master_rel.set("Type", notes_master_rel_type)
                    notes_master_rel.set("Target", "notesMasters/notesMaster1.xml")
                    rels_changed = True

                if rels_changed:
                    rels_tree.write(
                        rels_path,
                        xml_declaration=True,
                        encoding="UTF-8",
                        standalone="yes",
                    )

                remove_printer_settings_parts(ppt_dir, printer_settings_targets)
                remove_printer_settings_content_type(extract_dir)

                presentation_tree = etree.parse(presentation_path)
                presentation_root = presentation_tree.getroot()
                presentation_changed = False
                slide_size = presentation_root.find(f"{{{presentation_ns}}}sldSz")
                if slide_size is not None and slide_size.get("type") is not None:
                    del slide_size.attrib["type"]
                    presentation_changed = True

                if has_notes_master and notes_master_rel is not None:
                    notes_master_id_lst = presentation_root.find(
                        f"{{{presentation_ns}}}notesMasterIdLst"
                    )
                    if notes_master_id_lst is None:
                        notes_master_id_lst = etree.Element(
                            f"{{{presentation_ns}}}notesMasterIdLst"
                        )
                        sld_master_id_lst = presentation_root.find(
                            f"{{{presentation_ns}}}sldMasterIdLst"
                        )
                        if sld_master_id_lst is not None:
                            insert_index = list(presentation_root).index(
                                sld_master_id_lst
                            ) + 1
                            presentation_root.insert(insert_index, notes_master_id_lst)
                        else:
                            presentation_root.insert(0, notes_master_id_lst)
                    if not notes_master_id_lst.findall(
                        f"{{{presentation_ns}}}notesMasterId"
                    ):
                        notes_master_id = etree.SubElement(
                            notes_master_id_lst,
                            f"{{{presentation_ns}}}notesMasterId",
                        )
                        notes_master_id.set(
                            f"{{{rel_ns}}}id",
                            notes_master_rel.get("Id"),
                        )
                        presentation_changed = True

                if presentation_changed:
                    presentation_tree.write(
                        presentation_path,
                        xml_declaration=True,
                        encoding="UTF-8",
                        standalone="yes",
                    )

                normalize_presentation_relationship_ids(rels_path, presentation_path)

            ensure_slide_paragraphs_end_rpr(ppt_dir)
            normalize_zero_shadows_in_pptx(extract_dir)
            normalize_incomplete_themes(ppt_dir)

            with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as new_zip:
                for root, _, files in os.walk(extract_dir):
                    for file_name in files:
                        full_path = os.path.join(root, file_name)
                        archive_name = os.path.relpath(full_path, extract_dir)
                        new_zip.write(full_path, archive_name)

    async def fetch_network_assets(self):
        image_urls: List[str] = []
        models_with_network_asset: List[PptxPictureBoxModel] = []

        def process_image_path(shape_model: PptxPictureBoxModel):
            image_path = shape_model.picture.path

            local_path = resolve_image_path_to_filesystem(image_path)
            if local_path:
                shape_model.picture.path = local_path
                shape_model.picture.is_network = False
                return

            if is_data_image_uri(image_path):
                data_uri_path = write_data_image_to_file(image_path, self._temp_dir)
                if data_uri_path:
                    shape_model.picture.path = data_uri_path
                    shape_model.picture.is_network = False
                return

            if is_http_url(image_path):
                image_urls.append(image_path)
                models_with_network_asset.append(shape_model)

        if self._ppt_model.shapes:
            for shape_model in self._ppt_model.shapes:
                if isinstance(shape_model, PptxPictureBoxModel):
                    process_image_path(shape_model)

        for slide_model in self._slide_models:
            for shape_model in slide_model.shapes:
                if isinstance(shape_model, PptxPictureBoxModel):
                    process_image_path(shape_model)

        if image_urls:
            image_paths = await download_files(image_urls, self._temp_dir)
            for shape_model, image_path in zip(models_with_network_asset, image_paths):
                if image_path:
                    shape_model.picture.path = image_path
                    shape_model.picture.is_network = False

    async def create_ppt(self):
        await self.fetch_network_assets()

        global_shapes = list(self._ppt_model.shapes or [])
        for slide_model in self._slide_models:
            combined_slide_model = PptxSlideModel(
                background=slide_model.background,
                note=slide_model.note,
                shapes=[*slide_model.shapes, *global_shapes],
            )
            self.add_and_populate_slide(combined_slide_model)

    def add_and_populate_slide(self, slide_model: PptxSlideModel):
        slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[BLANK_SLIDE_LAYOUT])

        if slide_model.background:
            self.apply_fill_to_shape(slide.background, slide_model.background)

        if slide_model.note:
            slide.notes_slide.notes_text_frame.text = slide_model.note

        for shape_model in slide_model.shapes:
            if isinstance(shape_model, PptxPictureBoxModel):
                self.add_picture(slide, shape_model)
            elif isinstance(shape_model, PptxAutoShapeBoxModel):
                self.add_autoshape(slide, shape_model)
            elif isinstance(shape_model, PptxTextBoxModel):
                self.add_textbox(slide, shape_model)
            elif isinstance(shape_model, PptxConnectorModel):
                self.add_connector(slide, shape_model)

    def add_connector(self, slide: Slide, connector_model: PptxConnectorModel):
        if connector_model.thickness == 0:
            return
        connector_shape = slide.shapes.add_connector(
            connector_model.type,
            *connector_model.position.to_pt_xyxy(),
        )
        connector_shape.line.width = Pt(connector_model.thickness)
        connector_shape.line.color.rgb = RGBColor.from_string(connector_model.color)
        self.set_fill_opacity(connector_shape.line.fill, connector_model.opacity)

    def add_picture(self, slide: Slide, picture_model: PptxPictureBoxModel):
        image_path = picture_model.picture.path
        local_path = resolve_image_path_to_filesystem(image_path)
        if local_path:
            image_path = local_path
        elif is_data_image_uri(image_path):
            data_uri_path = write_data_image_to_file(image_path, self._temp_dir)
            if not data_uri_path:
                return
            image_path = data_uri_path

        if (
            picture_model.clip
            or picture_model.border_radius
            or picture_model.invert
            or picture_model.opacity
            or picture_model.object_fit
            or picture_model.shape
        ):
            try:
                image = Image.open(image_path)
            except Exception:
                return

            image = image.convert("RGBA")
            if picture_model.border_radius:
                image = round_image_corners(image, picture_model.border_radius)
            if picture_model.object_fit:
                image = fit_image(
                    image,
                    picture_model.position.width,
                    picture_model.position.height,
                    picture_model.object_fit,
                )
            elif picture_model.clip:
                image = clip_image(
                    image,
                    picture_model.position.width,
                    picture_model.position.height,
                )
            if picture_model.border_radius:
                image = round_image_corners(image, picture_model.border_radius)
            if picture_model.shape == PptxBoxShapeEnum.CIRCLE:
                image = create_circle_image(image)
            if picture_model.invert:
                image = invert_image(image)
            if picture_model.opacity is not None:
                image = set_image_opacity(image, picture_model.opacity)

            image_path = os.path.join(self._temp_dir, f"{uuid.uuid4()}.png")
            image.save(image_path)

        margined_position = self.get_margined_position(
            picture_model.position,
            picture_model.margin,
        )
        slide.shapes.add_picture(image_path, *margined_position.to_pt_list())

    def add_autoshape(self, slide: Slide, autoshape_box_model: PptxAutoShapeBoxModel):
        position = autoshape_box_model.position
        if autoshape_box_model.margin:
            position = self.get_margined_position(position, autoshape_box_model.margin)

        autoshape = slide.shapes.add_shape(
            autoshape_box_model.type,
            *position.to_pt_list(),
        )

        textbox = autoshape.text_frame
        textbox.word_wrap = autoshape_box_model.text_wrap

        self.apply_fill_to_shape(autoshape, autoshape_box_model.fill)
        self.apply_margin_to_text_box(textbox, autoshape_box_model.margin)
        self.apply_vertical_alignment_to_text_box(
            textbox, autoshape_box_model.vertical_alignment
        )
        self.apply_stroke_to_shape(autoshape, autoshape_box_model.stroke)
        self.apply_shadow_to_shape(autoshape, autoshape_box_model.shadow)
        self.apply_border_radius_to_shape(autoshape, autoshape_box_model.border_radius)

        if autoshape_box_model.paragraphs:
            self.add_paragraphs(textbox, autoshape_box_model.paragraphs)

    def add_textbox(self, slide: Slide, textbox_model: PptxTextBoxModel):
        position = textbox_model.position
        textbox_shape = slide.shapes.add_textbox(*position.to_pt_list())
        textbox_shape.width += Pt(2)

        textbox = textbox_shape.text_frame
        textbox.word_wrap = textbox_model.text_wrap

        self.apply_fill_to_shape(textbox_shape, textbox_model.fill)
        self.apply_margin_to_text_box(textbox, textbox_model.margin)
        self.apply_vertical_alignment_to_text_box(
            textbox, textbox_model.vertical_alignment
        )
        self.add_paragraphs(textbox, textbox_model.paragraphs)

    def add_paragraphs(
        self,
        textbox: TextFrame,
        paragraph_models: List[PptxParagraphModel],
    ):
        for index, paragraph_model in enumerate(paragraph_models):
            paragraph = textbox.add_paragraph() if index > 0 else textbox.paragraphs[0]
            self.populate_paragraph(paragraph, paragraph_model)

    def populate_paragraph(
        self,
        paragraph: _Paragraph,
        paragraph_model: PptxParagraphModel,
    ):
        if paragraph_model.spacing:
            self.apply_spacing_to_paragraph(paragraph, paragraph_model.spacing)

        if paragraph_model.line_height:
            paragraph.line_spacing = paragraph_model.line_height

        if paragraph_model.alignment:
            paragraph.alignment = paragraph_model.alignment

        if paragraph_model.font:
            self.apply_font_to_paragraph(paragraph, paragraph_model.font)

        text_runs: List[PptxTextRunModel] = []
        if paragraph_model.text:
            text_runs = self.parse_html_text_to_text_runs(
                paragraph_model.text,
                paragraph_model.font,
            )
        elif paragraph_model.text_runs:
            text_runs = paragraph_model.text_runs

        for text_run_model in text_runs:
            text_run = paragraph.add_run()
            self.populate_text_run(text_run, text_run_model)

    def parse_html_text_to_text_runs(
        self,
        text: str,
        font: Optional[PptxFontModel] = None,
    ) -> List[PptxTextRunModel]:
        return parse_html_text_to_text_runs(text, font)

    def populate_text_run(self, text_run: _Run, text_run_model: PptxTextRunModel):
        text_run.text = text_run_model.text
        if text_run_model.font:
            self.apply_font(text_run.font, text_run_model.font)

    def apply_border_radius_to_shape(self, shape: Shape, border_radius: Optional[int]):
        if not border_radius:
            return
        try:
            normalized_border_radius = Pt(border_radius) / min(shape.width, shape.height)
            shape.adjustments[0] = normalized_border_radius
        except Exception:
            return

    def apply_fill_to_shape(self, shape: Shape, fill: Optional[PptxFillModel] = None):
        if not fill:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.color)
            self.set_fill_opacity(shape.fill, fill.opacity)

    def apply_stroke_to_shape(
        self,
        shape: Shape,
        stroke: Optional[PptxStrokeModel] = None,
    ):
        if not stroke or stroke.thickness == 0:
            shape.line.fill.background()
        else:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = RGBColor.from_string(stroke.color)
            shape.line.width = Pt(stroke.thickness)
            self.set_fill_opacity(shape.line.fill, stroke.opacity)

    def apply_shadow_to_shape(
        self,
        shape: Shape,
        shadow: Optional[PptxShadowModel] = None,
    ):
        sp_element = shape._element
        sp_pr = sp_element.xpath("p:spPr")[0]
        nsmap = sp_pr.nsmap

        effect_list = sp_pr.find("a:effectLst", namespaces=nsmap)
        if effect_list:
            for node_name in ("a:outerShdw", "a:innerShdw", "a:prstShdw"):
                old_node = effect_list.find(node_name, namespaces=nsmap)
                if old_node is not None:
                    effect_list.remove(old_node)

        if not effect_list:
            effect_list = etree.SubElement(
                sp_pr,
                f"{{{nsmap['a']}}}effectLst",
                nsmap=nsmap,
            )

        if shadow is None:
            outer_shadow = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                nsmap=nsmap,
            )
            color_element = etree.SubElement(
                outer_shadow,
                f"{{{nsmap['a']}}}srgbClr",
                {"val": "000000"},
                nsmap=nsmap,
            )
            etree.SubElement(
                color_element,
                f"{{{nsmap['a']}}}alpha",
                {"val": "0"},
                nsmap=nsmap,
            )
            return

        angle_dir = (
            int(round((shadow.angle % 360) * 60000))
            if shadow.angle is not None
            else 0
        )
        outer_shadow = etree.SubElement(
            effect_list,
            f"{{{nsmap['a']}}}outerShdw",
            {
                "blurRad": f"{Pt(shadow.radius)}",
                "dir": f"{angle_dir}",
                "dist": f"{Pt(shadow.offset)}",
                "rotWithShape": "0",
            },
            nsmap=nsmap,
        )
        color_element = etree.SubElement(
            outer_shadow,
            f"{{{nsmap['a']}}}srgbClr",
            {"val": f"{shadow.color}"},
            nsmap=nsmap,
        )
        etree.SubElement(
            color_element,
            f"{{{nsmap['a']}}}alpha",
            {"val": f"{int(shadow.opacity * 100000)}"},
            nsmap=nsmap,
        )

    def set_fill_opacity(self, fill, opacity):
        if opacity is None or opacity >= 1.0:
            return

        alpha = int(opacity * 100000)

        try:
            solid_fill = fill._xPr.solidFill
            srgb = solid_fill.get_or_change_to_srgbClr()
            self.get_sub_element(srgb, "a:alpha", val=str(alpha))
        except Exception:
            return

    def get_margined_position(
        self,
        position: PptxPositionModel,
        margin: Optional[PptxSpacingModel],
    ) -> PptxPositionModel:
        if not margin:
            return position

        left = position.left + margin.left
        top = position.top + margin.top
        width = max(position.width - margin.left - margin.right, 0)
        height = max(position.height - margin.top - margin.bottom, 0)

        return PptxPositionModel(left=left, top=top, width=width, height=height)

    def apply_margin_to_text_box(
        self,
        text_frame: TextFrame,
        margin: Optional[PptxSpacingModel],
    ):
        text_frame.margin_left = Pt(margin.left if margin else 0)
        text_frame.margin_right = Pt(margin.right if margin else 0)
        text_frame.margin_top = Pt(margin.top if margin else 0)
        text_frame.margin_bottom = Pt(margin.bottom if margin else 0)

    def apply_vertical_alignment_to_text_box(
        self,
        text_frame: TextFrame,
        vertical_alignment,
    ):
        if vertical_alignment is None:
            return

        text_frame.vertical_anchor = vertical_alignment

    def apply_spacing_to_paragraph(
        self,
        paragraph: _Paragraph,
        spacing: PptxSpacingModel,
    ):
        paragraph.space_before = Pt(spacing.top)
        paragraph.space_after = Pt(spacing.bottom)

    def apply_font_to_paragraph(self, paragraph: _Paragraph, font: PptxFontModel):
        self.apply_font(paragraph.font, font)

    def apply_font(self, font: Font, font_model: PptxFontModel):
        font.name = font_model.name
        font.color.rgb = RGBColor.from_string(font_model.color)
        font.italic = font_model.italic
        font.size = Pt(font_model.size)
        font.bold = (
            font_model.font_weight is not None and font_model.font_weight >= 600
        )
        if font_model.underline is not None:
            font.underline = bool(font_model.underline)
        if font_model.strike is not None:
            self.apply_strike_to_font(font, font_model.strike)

    def apply_strike_to_font(self, font: Font, strike: Optional[bool]):
        try:
            run_properties = font._element
            if strike is True:
                run_properties.set("strike", "sngStrike")
            elif strike is False:
                run_properties.set("strike", "noStrike")
        except Exception:
            return

    def save(self, path: str):
        self._ppt.save(path)
        self.fix_keynote_compatibility(path)
